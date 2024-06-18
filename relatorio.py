from flask import Response
from pptx import Presentation
from io import BytesIO

def criar_slide(res, name):
    doc = Presentation()

    slides = []

    slides.append(doc.slides.add_slide(doc.slide_layouts[0]))
    slides[0].placeholders[0].text = 'Relatório - Brigadeiro Perfeito'
    slides[0].placeholders[1].text = 'Relatório da predição'

    slides.append(doc.slides.add_slide(doc.slide_layouts[2]))
    slides[1].placeholders[0].text = f'O valor previsto para {name} foi "{res}"'

    #doc.save(f"{nome}.pptx")
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    response = Response(output.read())
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    response.headers['Content-Disposition'] = 'attachment; filename=relatorio.pptx'

    return response